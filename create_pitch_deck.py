from pptx import Presentation
from pptx.util import Inches
import os

# Define paths
ARTIFACTS_DIR = "/Users/pop7/.gemini/antigravity/brain/97f27e07-9e29-4a02-bfb9-3adeb1b61dd2"
OUTPUT_FILE = os.path.join(ARTIFACTS_DIR, "ThaiScamBench_Pitch_Deck.pptx")

# Map of slide order to image filenames
SLIDES = [
    "pitch_deck_cover_1765127203434.png",
    "problem_slide_1765128489922.png",
    "solution_slide_1765128516956.png",
    "app_integration_flow_1765129436299.png", # New Integration Flow Slide
    "usecase_slide_1765128773649.png",
    "technology_slide_1765128540184.png",
    "market_opportunity_slide_1765127243054.png",
    "business_model_slide_1765127266916.png",
    "traction_slide_1765128586082.png",
    "competition_slide_1765128675572.png",
    "roadmap_slide_1765128608881.png",
    "ask_slide_1765128630164.png",
    "closing_slide_1765128695748.png"
]

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for image_file in SLIDES:
        image_path = os.path.join(ARTIFACTS_DIR, image_file)
        if os.path.exists(image_path):
            # Create a blank slide
            slide_layout = prs.slide_layouts[6] # 6 is blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image covering the entire slide
            slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            print(f"Added slide from {image_file}")
        else:
            print(f"Warning: Image not found {image_path}")

    # Add Team placeholder slide (since we don't have an image for it)
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "👥 TEAM"
    content.text = "Built by Security & AI Experts\n\nFounders:\n- CEO: 10+ years fintech, ex-bank executive\n- CTO: AI/ML specialist, ex-Google/Meta\n- CPO: Product leader, cybersecurity background\n\nAdvisory Board:\n- Former Bank of Thailand executive\n- Cybersecurity expert\n- Fintech founder"

    # Save
    prs.save(OUTPUT_FILE)
    print(f"Successfully saved presentation to {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
